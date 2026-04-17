from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_FILE = PROJECT_ROOT / "Vision_Based_Speed_Estimation_Traffic_Analysis_Presentation_v2.pptx"
RESULT_IMAGE = next((PROJECT_ROOT / "outputs").glob("speed_distribution_*.png"), None)

BG = RGBColor(20, 27, 32)
PANEL = RGBColor(27, 35, 43)
ACCENT = RGBColor(0, 166, 166)
ACCENT_2 = RGBColor(64, 201, 180)
TEXT = RGBColor(241, 246, 248)
MUTED = RGBColor(191, 204, 214)


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.42))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(11.8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.25), Inches(11.3), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(12)
        r.font.color.rgb = ACCENT_2


def add_bullets(slide, items: list[str], x=0.9, y=1.8, w=11.4, h=4.8, font_size=21):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)


def add_panel(slide, x, y, w, h, title, body_lines):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.3)

    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.3), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2

    body = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.5), Inches(w - 0.3), Inches(h - 0.6))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)


def add_tech_panel(slide, x, y, w, h, title, subtitle):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.1)

    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.16), Inches(w - 0.3), Inches(0.34))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2

    sub_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.56), Inches(w - 0.3), Inches(h - 0.7))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.size = Pt(13)
    r.font.color.rgb = TEXT


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    accent_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.3), Inches(2.55), Inches(0.45))
    accent_box.fill.solid()
    accent_box.fill.fore_color.rgb = PANEL
    accent_box.line.color.rgb = ACCENT
    t = accent_box.text_frame.paragraphs[0]
    t.alignment = PP_ALIGN.CENTER
    r = t.add_run()
    r.text = "Machine Vision Lab"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2
    add_title(slide, "Vision-Based Speed Estimation & Traffic Analysis", "Prepared from the finalized project implementation and report")
    add_bullets(
        slide,
        [
            "Submitted for the Machine Vision mini project",
            "Core pipeline: YOLOv8 detection, tracker-based IDs, homography speed estimation, live web analytics",
            "Focus areas: vehicle speed estimation, overspeeding alerts, wrong-way detection, illegal U-turn analysis, traffic counting",
        ],
        y=2.15,
        h=2.3,
        font_size=20,
    )
    add_panel(slide, 0.8, 5.2, 4.0, 1.2, "Team", ["Sameer Mishra (23BAI0221)", "Bhavyendra Manish (23BAI0176)"])
    add_panel(slide, 4.95, 5.2, 3.4, 1.2, "Department", ["Computer Science & Engineering", "Vellore Institute of Technology"])
    add_panel(slide, 8.55, 5.2, 3.9, 1.2, "Deliverable", ["Web-based traffic analysis dashboard", "Video upload and live camera support"])

    slides = [
        (
            "Problem Statement",
            [
                "Traffic monitoring requires estimating vehicle speed and traffic flow from ordinary camera footage without using physical road sensors.",
                "The system must detect vehicles, track them across frames, convert image motion into real-world speed, and present traffic analytics in real time.",
                "It must remain useful under changing lighting, different camera angles, partial occlusion, and mixed vehicle sizes.",
            ],
        ),
        (
            "Literature Review",
            [
                "Optical-flow methods estimate motion vectors directly from pixel movement but become noisy under lighting variation and cluttered scenes.",
                "Background subtraction and frame differencing are lightweight and fast, but they fail with dynamic backgrounds and camera instability.",
                "Recent work combines deep detectors such as YOLO with tracking methods like SORT or ByteTrack for stronger vehicle detection and ID consistency.",
                "Stereo and LiDAR-assisted systems improve distance accuracy, but they need specialized hardware and higher deployment cost.",
            ],
        ),
        (
            "Gap Identified",
            [
                "Many published systems depend on precise camera calibration and known road reference points, which is difficult for general deployment.",
                "Performance drops in rain, fog, shadows, and dense traffic where trackers can lose vehicle identity.",
                "Several solutions stop at speed estimation and do not expose broader traffic analytics such as counting, density, congestion, and live dashboard output.",
                "Real-time performance on standard hardware remains a challenge for heavy deep-learning pipelines.",
            ],
        ),
        (
            "Objective",
            [
                "Detect vehicles in traffic footage using a YOLOv8-based detector.",
                "Track each vehicle consistently across frames using tracker-based persistent IDs.",
                "Estimate speed in km/h through perspective transformation and displacement over time.",
                "Count vehicles crossing a virtual line and display live traffic analytics on a web dashboard.",
                "Flag overspeeding, wrong-way movement, and possible illegal U-turns with confidence values.",
            ],
        ),
        (
            "Tech Stack",
            [],
        ),
        (
            "Datasets",
            [
                "Primary traffic video dataset: UA-DETRAC, a public benchmark with 100 clips and more than 140,000 frames under clear, rainy, overcast, and night conditions.",
                "Supplementary footage: recorded traffic videos used to test the live dashboard and end-to-end workflow.",
                "Violation-model recalibration: Indian_Traffic_Violations.csv used to strengthen overspeed-related confidence and speed-threshold behavior.",
                "Track-feature seed data is retained for wrong-way and U-turn decision patterns where the Indian CSV does not provide explicit labels.",
            ],
        ),
        (
            "Proposed Methodology Workflow",
            [
                "1. Accept a traffic video file or live camera stream as input.",
                "2. Define a Region of Interest (ROI) and read frames using OpenCV.",
                "3. Run YOLOv8 to detect vehicles and ByteTrack persistence to maintain IDs.",
                "4. Store vehicle centres over time and project ground points using a homography matrix.",
                "5. Estimate speed from transformed displacement over time, then compare against the speed limit.",
                "6. Compute traffic analytics: active vehicles, total count, average speed, congestion, heatmap, alerts, and confidence levels.",
            ],
        ),
        (
            "Implementation (Demo)",
            [
                "Backend: Python, Flask, OpenCV, NumPy, YOLOv8 (Ultralytics), tracker-based ID persistence, scikit-learn, and Matplotlib.",
                "Frontend: web dashboard with source selection, calibration inputs, live annotated video stream, metrics cards, and alert panel.",
                "Video overlays include bounding boxes, IDs, speed labels, confidence labels, ROI guide, counting line, and heatmap.",
                "Stopping a session automatically generates a speed-distribution graph for report/demo support.",
            ],
        ),
        (
            "Result",
            [
                "The finalized project now matches the report pipeline more closely: YOLOv8 detection, tracker-based IDs, homography speed estimation, counting line, heatmap, and dashboard analytics are implemented.",
                "The web dashboard reports active vehicles, total vehicle count, average speed, congestion level, violation counts, overall confidence, and alert confidence.",
                "A speed distribution plot is generated for completed sessions, supporting result visualization beyond the live dashboard.",
                "Performance and exact benchmark numbers depend on the test clip, hardware, and camera calibration quality.",
            ],
        ),
        (
            "Conclusion",
            [
                "A practical camera-based traffic monitoring pipeline can estimate vehicle speed and deliver useful analytics without physical road sensors.",
                "The final system combines modern vehicle detection, tracking, homography-based measurement, and a web dashboard for real-time monitoring.",
                "The strongest benefits are cost-effective deployment, reusable CCTV infrastructure, and flexible traffic analysis from standard video input.",
                "Future improvements can include lane-wise analytics, stronger benchmark evaluation on UA-DETRAC, and automated calibration for improved speed accuracy.",
            ],
        ),
    ]

    for title, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)
        add_title(slide, title)
        if title == "Tech Stack":
            add_title(slide, title, "Core technologies used to build the finalized system")
            add_tech_panel(slide, 0.8, 1.8, 2.45, 1.3, "Python", "Backend logic, analytics, model integration, and orchestration")
            add_tech_panel(slide, 3.45, 1.8, 2.45, 1.3, "Flask", "Web server, API routes, live video stream, and dashboard integration")
            add_tech_panel(slide, 6.1, 1.8, 2.45, 1.3, "OpenCV", "Frame processing, homography mapping, overlays, counting line, and heatmap rendering")
            add_tech_panel(slide, 8.75, 1.8, 3.0, 1.3, "YOLOv8 + ByteTrack", "Vehicle detection and persistent multi-object tracking with stable IDs")
            add_tech_panel(slide, 0.8, 3.45, 2.45, 1.3, "NumPy", "Numerical computation for displacement, speed, and matrix operations")
            add_tech_panel(slide, 3.45, 3.45, 2.45, 1.3, "scikit-learn", "Violation confidence model and decision support for suspicious behavior")
            add_tech_panel(slide, 6.1, 3.45, 2.45, 1.3, "Matplotlib", "Speed distribution plot generation for result visualization")
            add_tech_panel(slide, 8.75, 3.45, 3.0, 1.3, "HTML / CSS / JavaScript", "Frontend dashboard, controls, metrics cards, and live alerts")
            add_panel(
                slide,
                0.82,
                5.35,
                10.93,
                1.05,
                "Integration Summary",
                ["The project combines computer vision, tracking, calibration, machine learning, and a live web dashboard into one end-to-end traffic analysis system."],
            )
        else:
            add_bullets(slide, bullets)

    # Enhance workflow slide with a visual flow row
    workflow_slide = prs.slides[7]
    steps = ["Input", "Detection", "Tracking", "Homography", "Speed", "Dashboard"]
    start_x = 0.95
    for i, step in enumerate(steps):
        x = start_x + i * 1.95
        box = workflow_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.5), Inches(1.5), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = ACCENT_2
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = step
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT
        if i < len(steps) - 1:
            arrow = workflow_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 1.57), Inches(5.64), Inches(0.28), Inches(0.36))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # Results slide add chart if available
    result_slide = prs.slides[10]
    if RESULT_IMAGE and RESULT_IMAGE.exists():
        result_slide.shapes.add_picture(str(RESULT_IMAGE), Inches(8.2), Inches(1.75), width=Inches(4.2), height=Inches(2.7))
        add_panel(result_slide, 8.15, 4.7, 4.25, 1.15, "Output Artifact", ["Automatically generated speed-distribution chart", "Created when an analysis session ends"])

    prs.save(OUTPUT_FILE)
    print(OUTPUT_FILE)


if __name__ == "__main__":
    build()
